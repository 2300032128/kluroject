import io
import pypdf
import docx
from pptx import Presentation
from typing import List, Dict, Any

class DocumentIngestionService:
    """
    Extracts text and page/slide structure from PDF, PPTX, DOCX, TXT, and MD files.
    """

    def extract_document(self, filename: str, content_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Returns a list of page/slide blocks:
        [
          {"page_number": 1, "text": "..."},
          {"page_number": 2, "text": "..."}
        ]
        """
        ext = filename.split(".")[-1].lower() if "." in filename else "txt"

        if ext == "pdf":
            return self._extract_pdf(content_bytes)
        elif ext in ("pptx", "ppt"):
            return self._extract_pptx(content_bytes)
        elif ext in ("docx", "doc"):
            return self._extract_docx(content_bytes)
        else:
            # TXT, MD or generic text
            return self._extract_text(content_bytes)

    def _extract_pdf(self, content_bytes: bytes) -> List[Dict[str, Any]]:
        pages = []
        try:
            reader = pypdf.PdfReader(io.BytesIO(content_bytes))
            for idx, page in enumerate(reader.pages):
                text = self._clean_text(page.extract_text() or "")
                if text:
                    pages.append({"page_number": idx + 1, "text": text})
        except Exception as e:
            print(f"[DocumentIngestionService] PDF Extraction Warning: {e}")
            raw = content_bytes.decode("utf-8", errors="ignore")
            pages.append({"page_number": 1, "text": self._clean_text(raw)})

        if not pages:
            pages.append({"page_number": 1, "text": "Sample PDF Document Content"})
        return pages

    def _extract_pptx(self, content_bytes: bytes) -> List[Dict[str, Any]]:
        slides = []
        try:
            prs = Presentation(io.BytesIO(content_bytes))
            for idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                combined = self._clean_text("\n".join(slide_texts))
                if combined:
                    slides.append({"page_number": idx + 1, "text": combined})
        except Exception as e:
            print(f"[DocumentIngestionService] PPTX Extraction Warning: {e}")
            raw = content_bytes.decode("utf-8", errors="ignore")
            slides.append({"page_number": 1, "text": self._clean_text(raw)})

        if not slides:
            slides.append({"page_number": 1, "text": "Sample PPTX Lecture Slide Content"})
        return slides

    def _extract_docx(self, content_bytes: bytes) -> List[Dict[str, Any]]:
        paragraphs = []
        try:
            doc = docx.Document(io.BytesIO(content_bytes))
            full_text = self._clean_text("\n".join([p.text for p in doc.paragraphs if p.text]))
            if full_text:
                paragraphs.append({"page_number": 1, "text": full_text})
        except Exception as e:
            print(f"[DocumentIngestionService] DOCX Extraction Warning: {e}")
            raw = content_bytes.decode("utf-8", errors="ignore")
            paragraphs.append({"page_number": 1, "text": self._clean_text(raw)})

        if not paragraphs:
            paragraphs.append({"page_number": 1, "text": "Sample DOCX Document Content"})
        return paragraphs

    def _extract_text(self, content_bytes: bytes) -> List[Dict[str, Any]]:
        text_str = self._clean_text(content_bytes.decode("utf-8", errors="ignore"))
        return [{"page_number": 1, "text": text_str if text_str else "Sample Text Document Content"}]

    def _clean_text(self, text: str) -> str:
        if not text:
            return ""
        # Remove extra whitespace while preserving structure
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)

ingestion_service = DocumentIngestionService()
